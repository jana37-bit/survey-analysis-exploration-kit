#!/usr/bin/env python3
"""
Build a survey research presentation from analysis results.

Usage:
    build_presentation.py <results.json> <output.pptx> [options]

Options:
    --template <file.pptx>    Use existing template
    --title "Project Name"    Override title

The results.json should contain:
{
    "title": "Survey Analysis",
    "date": "February 2026",
    "n_respondents": 2227,
    "methodology": "Online survey, 18+ adults",
    "key_findings": [
        {"stat": "72%", "description": "satisfied with service"},
        ...
    ],
    "charts": [
        {"path": "charts/satisfaction.png", "title": "Satisfaction", "base": 500},
        ...
    ]
}

Examples:
    build_presentation.py outputs/results.json outputs/summary.pptx
    build_presentation.py outputs/results.json outputs/summary.pptx --title "Q4 Brand Tracker"
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Missing package. Run: pip install python-pptx")
    sys.exit(1)

# Colors
PRIMARY = RGBColor(0x2F, 0x54, 0x96)
SECONDARY = RGBColor(0x44, 0x72, 0xC4)
ACCENT = RGBColor(0xED, 0x7D, 0x31)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_TEXT = RGBColor(0x66, 0x66, 0x66)
SOURCE_TEXT = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)


def add_title_slide(prs, title, subtitle):
    """Dark blue title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PRIMARY

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = ICE_BLUE
    p2.alignment = PP_ALIGN.LEFT


def add_exec_summary(prs, findings):
    """Executive summary with stat callouts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Key Findings"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Stat callouts (up to 3 across)
    for i, finding in enumerate(findings[:3]):
        left = Inches(0.5 + i * 4.2)

        box = slide.shapes.add_textbox(left, Inches(2.2), Inches(3.8), Inches(3))
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = str(finding.get("stat", ""))
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = str(finding.get("description", ""))
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_TEXT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)


def add_divider_slide(prs, section_title):
    """Section divider with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PRIMARY

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    p = txBox.text_frame.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_chart_slide(prs, chart_info):
    """Data slide with embedded chart image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = chart_info.get("title", "")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Chart image
    chart_path = chart_info.get("path", "")
    if chart_path and Path(chart_path).exists():
        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.3), Inches(12), Inches(5.2)
        )

    # Insight text (optional)
    insight = chart_info.get("insight", "")
    if insight:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = insight
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = DARK_TEXT

    # Source line
    base = chart_info.get("base", "")
    source = chart_info.get("source", "Survey Analysis")
    source_text = f"Source: {source}"
    if base:
        source_text += f", n={base}"

    src_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
    p = src_box.text_frame.paragraphs[0]
    p.text = source_text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = SOURCE_TEXT


def build_presentation(results, output_path, template_path=None):
    """Build complete presentation from results dict."""
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    title = results.get("title", "Survey Analysis")
    date = results.get("date", "")
    n = results.get("n_respondents", "")
    subtitle = f"n={n}" + (f" | {date}" if date else "")

    add_title_slide(prs, title, subtitle)

    findings = results.get("key_findings", [])
    if findings:
        add_exec_summary(prs, findings)

    charts = results.get("charts", [])
    if charts:
        add_divider_slide(prs, "Detailed Findings")
        for chart in charts:
            add_chart_slide(prs, chart)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Presentation saved: {output_path} ({len(prs.slides)} slides)")


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    results_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    template_path = None
    custom_title = None

    args = sys.argv[3:]
    i = 0
    while i < len(args):
        if args[i] == "--template" and i + 1 < len(args):
            template_path = args[i + 1]
            i += 2
        elif args[i] == "--title" and i + 1 < len(args):
            custom_title = args[i + 1]
            i += 2
        else:
            i += 1

    with open(results_path) as f:
        results = json.load(f)

    if custom_title:
        results["title"] = custom_title

    build_presentation(results, str(output_path), template_path)


if __name__ == "__main__":
    main()
